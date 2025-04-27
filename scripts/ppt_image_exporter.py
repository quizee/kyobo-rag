import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 이미지가 있는 폴더 경로
SECTIONS_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'test_output', 'sections')
OUTPUT_PPTX = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'test_output', 'test-output.pptx')

def create_footer_text(slide, metadata, presentation):
    # 텍스트 상자 추가
    left = Inches(0.5)  # 왼쪽 여백
    top = presentation.slide_height - Inches(0.8)  # 하단에서 0.8인치 위
    width = presentation.slide_width - Inches(1)  # 슬라이드 너비 - 1인치
    height = Inches(0.5)  # 텍스트 상자 높이
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    
    # 헤더 텍스트 생성
    header_text = metadata.get("Header 1", "")
    if metadata.get("Header 2"):
        header_text += f" - {metadata['Header 2']}"
    if metadata.get("Header 3"):
        header_text += f" - {metadata['Header 3']}"
    
    # 첫 번째 단락 (헤더 텍스트)
    p1 = tf.paragraphs[0]
    p1.text = header_text
    p1.font.size = Pt(8)
    p1.font.color.rgb = RGBColor(128, 128, 128)
    
    # 두 번째 단락 (AI 생성 텍스트)
    p2 = tf.add_paragraph()
    p2.text = "이 슬라이드는 AI가 생성했습니다"
    p2.font.size = Pt(8)
    p2.font.color.rgb = RGBColor(128, 128, 128)

def find_matching_image(image_files, target_name):
    target_number = ''.join(filter(str.isdigit, target_name))  # "image1" -> "1"
    for img_file in image_files:
        file_number = ''.join(filter(str.isdigit, img_file.split('.')[0]))  # "imag1.png" -> "1"
        if target_number == file_number:
            return img_file
    
    return None

def create_ppt_with_metadata(image_dir, output_pptx, metadata_list):
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]  # 빈 슬라이드

    # 이미지 파일 목록
    image_files = [f for f in os.listdir(image_dir) if f.lower().endswith((".png", ".jpg", ".jpeg"))]

    for metadata in metadata_list:
        slide = prs.slides.add_slide(blank_slide_layout)
        image_name = metadata.get("image_name")
        
        # 이미지 파일 찾기
        if image_name:
            matching_image = find_matching_image(image_files, image_name)
            
            if matching_image:
                img_path = os.path.join(image_dir, matching_image)
                
                # 슬라이드 크기
                slide_width = prs.slide_width
                slide_height = prs.slide_height

                # 이미지 크기 자동 조정
                from PIL import Image
                with Image.open(img_path) as im:
                    img_width, img_height = im.size
                    slide_ratio = slide_width / slide_height
                    img_ratio = img_width / img_height

                    if img_ratio > slide_ratio:
                        width = slide_width * 0.85
                        height = width / img_ratio
                    else:
                        height = slide_height * 0.85
                        width = height * img_ratio

                left = int((slide_width - width) / 2)
                top = int((slide_height - height) / 2)
                slide.shapes.add_picture(img_path, left, top, int(width), int(height))

        # 푸터 텍스트 추가
        create_footer_text(slide, metadata, prs)

    prs.save(output_pptx)
    print(f"PPT 저장 완료: {output_pptx}")

if __name__ == "__main__":
    # 테스트용 메타데이터
    test_metadata = [
        {
            "Header 1": "특약 안내",
            "Header 2": "뇌심치료",
            "page_number": 7,
            "image_name": "image1",
        },
        {
            "Header 1": "특약 안내",
            "Header 2": "암진단",
            "page_number": 7,
            "image_name": "image2",
        },
        {
            "Header 1": "특약 안내",
            "Header 2": "특정치료",
            "page_number": 7,
            "image_name": "image3",
        },
        {
            "Header 1": "특약 안내",
            "Header 2": "특정치료",
            "page_number": 7,
            "image_name": "image5",
        }
    ]
    create_ppt_with_metadata(SECTIONS_DIR, OUTPUT_PPTX, test_metadata) 