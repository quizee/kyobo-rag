import streamlit as st
import os
from dotenv import load_dotenv
import nest_asyncio
import re
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from PIL import Image as PILImage
from pptx.enum.text import PP_ALIGN
import time
import pandas as pd


nest_asyncio.apply()
from llama_cloud_services import LlamaParse
from langchain_text_splitters import MarkdownHeaderTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_openai import OpenAIEmbeddings
import openai
import json
from scripts.clova_ocr_utils import ocr_image
import requests
from scripts.upstage_parser import UpstageParser

# .env 환경변수 로드
load_dotenv()

st.set_page_config(page_title="PDF Chat App", layout="wide")


def crop_image_by_y(
    image_path, y_start, y_end, output_path, margin_top=15, margin_side=30
):
    """이미지를 y 좌표 기준으로 자르는 함수"""
    try:
        img = Image.open(image_path)
        img_width, img_height = img.size

        # 여백 적용 (부동소수점 유지)
        y_start = max(0, y_start - margin_top)
        y_end = min(img_height, y_end)
        x_start = max(0, margin_side)
        x_end = min(img_width, img_width - margin_side)

        if y_end <= y_start or x_end <= x_start:
            print("Invalid crop dimensions")
            return None

        # 최종적으로 정수로 변환 (이미지 크롭 시에만)
        crop_box = (int(x_start), int(y_start), int(x_end), int(y_end))
        cropped = img.crop(crop_box)
        cropped.save(output_path)
        return output_path
    except Exception as e:
        print(f"Error cropping image: {str(e)}")
        return None


def is_useless_header(text):
    # 한글이 하나라도 있으면 False
    if re.search(r"[가-힣]", text):
        return False
    # 한글이 없고, 숫자/영어/공백/특수문자만 있으면 True
    if re.fullmatch(r"[\dA-Za-z\s\W]+", text):
        return True
    return False


def build_header_dictionary(headers, _, output_dir="temp_output"):
    """
    headers: [{page, level, text, y, bbox}, ...]
    return: {"header_text": {page, level, text, y, crop_image_path, Header 1, Header 2, Header 3}}
    """
    header_dict = {}
    from collections import defaultdict

    # 페이지별로 그룹화 (1페이지 제외)
    page_to_headers = defaultdict(list)
    for header in headers:
        if (
            header["y"] is not None and header["page"] > 1
        ):  # y 좌표가 있고 1페이지가 아닌 경우만 처리
            page_to_headers[header["page"]].append(header)

    # 각 페이지별로 헤더 정보 수집 (1페이지 제외)
    for page_number, page_headers_list in page_to_headers.items():
        if page_number <= 1:  # 1페이지는 건너뛰기
            continue

        # y좌표 기준으로 정렬
        sorted_headers = sorted(page_headers_list, key=lambda x: x["y"])

        # 현재 페이지의 헤더 계층 구조 추적
        current_headers = {1: "", 2: "", 3: ""}  # Header 1  # Header 2  # Header 3

        for header in sorted_headers:
            level = header["level"]
            text = header["text"]

            # 현재 레벨의 헤더 업데이트
            current_headers[level] = text

            # 상위 레벨의 헤더 정보 유지
            for l in range(level + 1, 4):
                current_headers[l] = ""

            # 헤더 정보 저장
            key = f"{text}"
            header_dict[key] = {
                "page": page_number,
                "level": level,
                "text": text,
                "markdown_text": header.get("markdown_text", f"{'#' * level} {text}"),
                "y": header["y"],
                "crop_image_path": None,  # 나중에 설정됨
                "Header 1": current_headers[1],
                "Header 2": current_headers[2],
                "Header 3": current_headers[3],
                "page_number": page_number,
                "bbox": header.get("bbox"),
            }

    # 각 페이지의 헤더별로 이미지 생성
    for page_number, page_headers_list in page_to_headers.items():
        page_img_path = os.path.join(output_dir, f"page_{page_number}.jpg")

        # 페이지 이미지가 없으면 PDF에서 변환
        if not os.path.exists(page_img_path):
            print(f"Page image not found: {page_img_path}, converting from PDF...")
            # PDF 파일 찾기
            pdf_files = [
                f
                for f in os.listdir(output_dir)
                if f.startswith("temp_") and f.endswith(".pdf")
            ]
            if not pdf_files:
                print("No PDF file found in output directory")
                continue

            pdf_path = os.path.join(output_dir, pdf_files[0])
            # 현재 페이지만 변환
            try:
                import fitz  # PyMuPDF

                doc = fitz.open(pdf_path)
                if page_number <= len(doc):
                    page = doc[page_number - 1]  # 0-based index
                    pix = page.get_pixmap(matrix=fitz.Matrix(300 / 72, 300 / 72))
                    pix.save(page_img_path)
                    print(f"Converted page {page_number} to image: {page_img_path}")
                doc.close()
            except Exception as e:
                print(f"Error converting page {page_number} to image: {str(e)}")
                continue

        if not os.path.exists(page_img_path):
            print(f"Warning: Page image not found after conversion: {page_img_path}")
            continue

        print(f"Processing page {page_number}, image path: {page_img_path}")
        img = Image.open(page_img_path)
        img_width, img_height = img.size
        print(f"Image size: {img.size}")

        # 모든 헤더를 y좌표 기준으로 정렬
        all_headers = sorted(page_headers_list, key=lambda h: h["y"])
        print(f"Total headers on page {page_number}: {len(all_headers)}")

        # 각 헤더에 대해 이미지 생성
        for i, header in enumerate(all_headers):
            # 현재 헤더의 y 좌표 (상대값을 픽셀값으로 변환)
            y_start = header["y"] * img_height  # 정수 변환하지 않고 부동소수점 유지

            # 다음 헤더의 y 좌표 찾기 (헤더 레벨 고려)
            y_end = img_height  # 기본값은 페이지 끝
            current_header_text = header["text"]
            current_header_level = header["level"]

            for next_header in all_headers[i + 1 :]:
                next_y = next_header["y"] * img_height
                if next_y > y_start:  # 현재 헤더보다 아래에 있는 헤더
                    if next_header["level"] == current_header_level:
                        y_end = next_y
                        break

            print(
                f"Header {i+1}: y_start={y_start:.2f}, y_end={y_end:.2f}, text={header['text'][:20]}..."
            )

            if y_end <= y_start:
                print(f"Skipping header {i+1}: invalid y coordinates")
                continue

            # 이미지 여백 계산 (이미지 크기에 비례)
            margin_top = img_height * 0.02  # 상단 여백 2%
            margin_side = img_width * 0.05  # 좌우 여백 5%

            # 바운딩 박스가 있는 경우, 좌우 여백을 바운딩 박스 기반으로 조정
            if header.get("bbox"):
                bbox = header["bbox"]
                try:
                    # 좌우 좌표를 픽셀값으로 변환
                    left_x = bbox["top_left"]["x"] * img_width
                    right_x = bbox["top_right"]["x"] * img_width
                    # 바운딩 박스 너비의 5%를 여백으로 사용
                    margin_side = (right_x - left_x) * 0.05
                except (KeyError, TypeError):
                    # bbox 구조가 예상과 다르거나 None인 경우 기본 여백 사용
                    margin_side = img_width * 0.05

            out_path = os.path.join(
                output_dir,
                f"page_{page_number}_level{header['level']}_header_{i+1}_{header['text'].replace('#', '').strip()}.jpg",
            )

            # 이미 크롭된 이미지가 있는지 확인
            if os.path.exists(out_path):
                print(f"Using existing cropped image: {out_path}")
            else:
                print(f"Cropping image to: {out_path}")
                crop_result = crop_image_by_y(
                    page_img_path,
                    y_start,
                    y_end,
                    out_path,
                    margin_top=margin_top,
                    margin_side=margin_side,
                )
                if crop_result is None:
                    print(f"Failed to crop image for header {i+1}")
                    continue
                print(f"Successfully cropped image: {out_path}")

            # 헤더 정보 업데이트
            header_dict[current_header_text]["crop_image_path"] = out_path

    print(f"Total headers processed: {len(header_dict)}")
    return header_dict


def add_source_text(slide, header_info, prs, pdf_name=""):
    """하단 중앙에 출처 텍스트 추가"""
    # page_number를 header_info에서 직접 가져오기 (page 또는 page_number 키 사용)
    page_number = header_info.get("page_number") or header_info.get("page", "?")

    # 출처 텍스트 박스
    width = Inches(6)  # 고정된 너비
    left = (prs.slide_width - width) / 2  # 중앙 정렬
    top = prs.slide_height - Inches(1.2)  # AI 고지 문구 위에 위치
    height = Inches(0.5)

    ref_box = slide.shapes.add_textbox(left, top, width, height)
    ref_tf = ref_box.text_frame
    ref_tf.word_wrap = True

    p_ref = ref_tf.paragraphs[0]
    p_ref.alignment = PP_ALIGN.CENTER

    run1 = p_ref.add_run()
    run1.text = "[출처: "
    run1.font.size = Pt(12)
    run1.font.color.rgb = RGBColor(120, 120, 120)

    run2 = p_ref.add_run()
    run2.text = f'"{pdf_name}", {page_number}page]'
    run2.font.size = Pt(12)
    run2.font.color.rgb = RGBColor(120, 120, 120)


def add_header_path(slide, header_info, prs):
    """제목 위에 계층적 헤더 경로 추가"""
    # level과 text를 사용해 계층적 헤더 구성
    level = header_info.get("level", 0)
    text = header_info.get("text", "")

    # 현재 헤더의 상위 헤더들을 찾기 위해 page와 y좌표 사용
    page = header_info.get("page", 0)
    y = header_info.get("y", 0)

    # 같은 페이지의 모든 헤더를 가져와서 y좌표 기준으로 정렬
    all_headers = []
    for key, info in st.session_state.get("header_dict", {}).items():
        if info.get("page") == page:
            all_headers.append(info)

    # y좌표 기준으로 정렬
    all_headers.sort(key=lambda x: x.get("y", 0))

    # 현재 헤더보다 위에 있는 헤더들 중 직계 상위 헤더만 찾기
    header_text = ""
    for h in all_headers:
        if h.get("y", 0) < y:
            h_level = h.get("level", 0)
            # 현재 헤더의 직계 상위 헤더만 추가
            if h_level == level - 1:
                if header_text:
                    header_text += " >> "
                header_text += h.get("text", "")

    # 현재 헤더 추가
    if header_text:
        header_text += " >> " + text
    else:
        header_text = text

    # 헤더 경로 텍스트 박스
    left = Inches(0.3)
    top = Inches(0.4)  # 제목보다 더 위에 위치
    width = prs.slide_width - Inches(0.6)
    height = Inches(0.5)

    path_box = slide.shapes.add_textbox(left, top, width, height)
    path_tf = path_box.text_frame
    path_tf.word_wrap = True

    p_path = path_tf.paragraphs[0]
    p_path.text = header_text
    p_path.font.size = Pt(12)  # 12pt로 변경
    p_path.font.bold = True
    p_path.font.color.rgb = RGBColor(0, 0, 255)  # 파란색
    p_path.alignment = PP_ALIGN.LEFT


def add_title_text(slide, header_info, prs):
    """상단에 큰 제목 추가"""
    # 계층적 헤더 경로에서 마지막 부분만 추출
    header_path = header_info.get("text", "")
    if " >> " in header_path:
        title_text = header_path.split(" >> ")[-1]
    else:
        title_text = header_path

    title_left = Inches(0.3)
    title_top = Inches(0.6)
    title_width = prs.slide_width - Inches(0.6)
    title_height = Inches(1)

    title_box = slide.shapes.add_textbox(
        title_left, title_top, title_width, title_height
    )
    title_tf = title_box.text_frame
    title_tf.word_wrap = True

    p_title = title_tf.paragraphs[0]
    p_title.text = title_text
    p_title.font.size = Pt(24)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(80, 80, 80)
    p_title.alignment = PP_ALIGN.LEFT


def add_ai_notice_text(slide, prs):
    """하단 중앙에 생성형 AI 고지 문구 추가"""
    bottom_text = (
        "본 자료는 생성형 AI 기반으로 작성되었으며, 중요한 사실은 확인이 필요합니다."
    )

    box_width = Inches(6)
    bottom_left = (prs.slide_width - box_width) / 2
    bottom_top = prs.slide_height - Inches(0.7)
    bottom_height = Inches(0.5)

    ai_box = slide.shapes.add_textbox(bottom_left, bottom_top, box_width, bottom_height)
    ai_tf = ai_box.text_frame
    ai_tf.word_wrap = True

    p_ai = ai_tf.paragraphs[0]
    p_ai.text = bottom_text
    p_ai.font.size = Pt(12)
    p_ai.font.color.rgb = RGBColor(150, 150, 150)
    p_ai.alignment = PP_ALIGN.CENTER


def add_center_image(slide, header_info, prs):
    """중앙에 이미지 추가"""
    img_path = header_info.get("crop_image_path")
    if img_path and os.path.exists(img_path):
        with PILImage.open(img_path) as im:
            img_width, img_height = im.size
            slide_width = prs.slide_width
            slide_height = prs.slide_height

            # 슬라이드의 40%를 이미지 최대 높이로 설정
            max_height = slide_height * 0.4

            # 이미지 비율 유지하면서 크기 조정
            img_ratio = img_width / img_height

            # 높이 기준으로 크기 계산
            height = max_height
            width = height * img_ratio

            # 너비가 슬라이드 너비의 85%를 넘으면 너비 기준으로 다시 계산
            if width > slide_width * 0.85:
                width = slide_width * 0.85
                height = width / img_ratio

            # kyobo_layout2 이미지 아래, 출처 텍스트 위에 위치하도록 설정
            # kyobo_layout2는 상단에서 Inches(1.2)에 위치
            # 출처 텍스트는 하단에서 Inches(1.2)에 위치
            layout_margin_top = Inches(1.2)  # kyobo_layout2의 상단 여백
            source_margin_bottom = Inches(1.2)  # 출처 텍스트의 하단 여백

            # 이미지를 kyobo_layout2와 출처 텍스트 사이의 중앙에 배치
            available_height = slide_height - layout_margin_top - source_margin_bottom
            top = layout_margin_top + (available_height - height) / 2

            # 가로 중앙 정렬
            left = (slide_width - width) / 2

            slide.shapes.add_picture(img_path, left, top, width, height)


def add_layout_image(slide, prs):
    """제목 아래에 레이아웃 이미지 추가"""
    layout_path = os.path.join("temp_output", "kyobo_layout.jpg")
    if os.path.exists(layout_path):
        with PILImage.open(layout_path) as im:
            img_width, img_height = im.size
            # 슬라이드 너비의 95%로 이미지 크기 조정
            width = prs.slide_width * 0.95
            height = (img_height / img_width) * width

            # 이미지를 가로 중앙에 배치하고, 제목 아래에 위치시킴
            left = (prs.slide_width - width) / 2
            top = Inches(1.2)  # 제목 아래 위치

            slide.shapes.add_picture(layout_path, left, top, width, height)


def add_muscle_image(slide, prs):
    """kyobo_muscle.jpg 이미지 추가"""
    muscle_path = os.path.join("temp_output", "kyobo_muscle.jpg")
    if os.path.exists(muscle_path):
        # cm를 EMU로 변환 (1cm = 360000 EMU)
        width_cm = 4.05
        height_cm = 0.97
        width = int(width_cm * 360000)
        height = int(height_cm * 360000)

        # 정확한 위치 지정 (가로 20.71cm, 세로 17.57cm)
        left = int(20.71 * 360000)  # 가로 위치
        top = int(17.57 * 360000)  # 세로 위치

        slide.shapes.add_picture(muscle_path, left, top, width, height)


def create_ppt_from_header_dict(header_dict, output_path):
    """header_dict를 기반으로 PPT 생성"""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    print(f"Creating PPT with {len(header_dict)} headers")

    # 제목 슬라이드 추가
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "교보마이플랜건강보험"
    subtitle.text = "문서 분석 결과"

    # 각 헤더별로 슬라이드 추가
    for header_text, header_info in header_dict.items():
        # Header 1(level==1)은 PPT에 포함하지 않음
        if header_info.get("level") == 1:
            continue
        print(f"\nProcessing header: {header_text[:30]}...")
        print(f"Header info: {json.dumps(header_info, indent=2, ensure_ascii=False)}")

        slide_layout = prs.slide_layouts[6]  # 빈 슬라이드 레이아웃 사용
        slide = prs.slides.add_slide(slide_layout)

        # 헤더 경로 추가
        add_header_path(slide, header_info, prs)

        # 제목 추가
        add_title_text(slide, header_info, prs)

        # 이미지 추가
        crop_image_path = header_info.get("crop_image_path")
        print(f"Checking image path: {crop_image_path}")

        if crop_image_path:
            if os.path.exists(crop_image_path):
                print(f"Image file exists: {crop_image_path}")
                try:
                    # 이미지 크기 확인
                    with Image.open(crop_image_path) as img:
                        img_width, img_height = img.size
                        print(f"Image size: {img_width}x{img_height}")

                        # 슬라이드 크기
                        slide_width = prs.slide_width
                        slide_height = prs.slide_height

                        # 이미지를 슬라이드 중앙에 배치
                        # 상단 여백 (제목 아래)
                        top_margin = Inches(2.0)
                        # 하단 여백 (출처 텍스트 위)
                        bottom_margin = Inches(1.5)
                        # 사용 가능한 높이
                        available_height = slide_height - top_margin - bottom_margin

                        # 이미지 비율 유지하면서 크기 조정
                        img_ratio = img_width / img_height

                        # 높이 기준으로 크기 계산
                        height = available_height
                        width = height * img_ratio

                        # 너비가 슬라이드 너비의 90%를 넘으면 너비 기준으로 다시 계산
                        if width > slide_width * 0.9:
                            width = slide_width * 0.9
                            height = width / img_ratio

                        # 이미지 위치 계산 (중앙 정렬)
                        left = (slide_width - width) / 2
                        top = top_margin + (available_height - height) / 2

                        print(
                            f"Adding image at: left={left}, top={top}, width={width}, height={height}"
                        )
                        slide.shapes.add_picture(
                            crop_image_path, left, top, width, height
                        )
                        print("Image added successfully")
                except Exception as e:
                    print(f"Error adding image: {str(e)}")
                    import traceback

                    print(traceback.format_exc())
            else:
                print(f"Image file does not exist: {crop_image_path}")
        else:
            print("No image path in header info")

        # 출처 텍스트 추가
        add_source_text(slide, header_info, prs, "교보마이플랜건강보험")

        # AI 고지 문구 추가
        add_ai_notice_text(slide, prs)

    # PPT 저장
    print(f"\nSaving PPT to: {output_path}")
    try:
        prs.save(output_path)
        print("PPT saved successfully")
    except Exception as e:
        print(f"Error saving PPT: {str(e)}")
        import traceback

        print(traceback.format_exc())

    return output_path


def get_deepest_header_and_level(metadata):
    """메타데이터에서 가장 깊은 레벨의 헤더를 찾되, 실제 헤더 레벨을 반환"""
    # 먼저 Header 1부터 Header 3까지 확인
    for depth in range(1, 4):  # 1부터 3까지 순차적으로 확인
        key = f"Header {depth}"
        if key in metadata and metadata[key]:
            # 실제 헤더 레벨을 찾기 위해 header_dict 확인
            header_text = metadata[key]
            for h in st.session_state.get("header_dict", {}).values():
                # 텍스트 비교 시 공백 제거하고 정규화
                if h["text"].strip() == header_text.strip():
                    return header_text, h["level"]  # 실제 헤더 레벨 반환
    return None, None


def convert_pdf_to_images(pdf_path, output_dir):
    """PDF 파일을 이미지로 변환하여 저장 (1페이지 제외)"""
    try:
        import fitz  # PyMuPDF

        print(f"Converting PDF to images: {pdf_path}")

        # PDF 열기
        doc = fitz.open(pdf_path)
        total_pages = len(doc)
        print(f"Total pages: {total_pages}")

        # 2페이지부터 시작 (1페이지 제외)
        for page_num in range(1, total_pages):  # 1부터 시작하여 1페이지 제외
            page = doc[page_num]
            # DPI를 300으로 설정하여 고품질 이미지 생성
            pix = page.get_pixmap(matrix=fitz.Matrix(300 / 72, 300 / 72))
            output_path = os.path.join(
                output_dir, f"page_{page_num + 1}.jpg"
            )  # page_num + 1로 파일명 유지
            pix.save(output_path)
            print(f"Saved page {page_num + 1} as image: {output_path}")

        doc.close()
        return True
    except Exception as e:
        print(f"Error converting PDF to images: {str(e)}")
        import traceback

        print(traceback.format_exc())
        return False


# 2단 컬럼 레이아웃 (왼쪽: 채팅, 오른쪽: 파일 업로드)
col_chat, col_upload = st.columns([2, 1])

# FAISS DB를 세션에 저장
if "faiss_db" not in st.session_state:
    st.session_state["faiss_db"] = None

with col_chat:
    st.header("💬 교육 자료 생성하기")
    if "chat_history" not in st.session_state:
        st.session_state["chat_history"] = []
    # 채팅 내역 표시
    for msg in st.session_state["chat_history"]:
        st.markdown(f"**{msg['role']}**: {msg['content']}")
    # 사용자 입력
    user_input = st.text_input("메시지를 입력하세요", key="user_input")

    # 검색 버튼 및 결과 출력
    if st.button("검색"):
        if st.session_state["faiss_db"] is not None and user_input:
            header_dict = st.session_state.get("header_dict", {})
            retriever = st.session_state["faiss_db"].as_retriever(
                search_type="similarity", search_kwargs={"k": 3}
            )
            results = retriever.invoke(user_input)
            st.subheader("검색 결과")
            for i, doc in enumerate(results):
                page = doc.metadata.get("page_number", "알 수 없음")
                st.markdown(f"**결과 {i+1} (페이지: {page})**")
                st.code(doc.page_content)
                # === 이미지 미리보기 추가 ===
                header_text, level = get_deepest_header_and_level(doc.metadata)
                matched_header = None
                for h in header_dict.values():
                    if (
                        h["page"] == page
                        and h["level"] == level
                        and h["text"] == header_text
                    ):
                        matched_header = h
                        break
                if (
                    matched_header
                    and matched_header.get("crop_image_path")
                    and os.path.exists(matched_header["crop_image_path"])
                ):
                    st.image(
                        matched_header["crop_image_path"],
                        caption=f"이미지 ({header_text})",
                        use_column_width=True,
                    )

            # === 검색 결과 PPT 다운로드 버튼 ===
            selected_headers = []

            for doc in results:
                page = doc.metadata.get("page_number")
                header_text, level = get_deepest_header_and_level(doc.metadata)
                for h in header_dict.values():
                    if (
                        h["page"] == page
                        and h["level"] == level
                        and h["text"].strip() == header_text.strip()
                    ):
                        selected_headers.append(h)
                        break

            st.session_state["selected_headers"] = selected_headers
            # 검색 시점에 PPT를 미리 생성하고 경로를 세션에 저장
            pptx_path = os.path.join("temp_output", "search_results.pptx")

            if selected_headers:
                # selected_headers를 딕셔너리 형태로 변환
                selected_headers_dict = {}
                for header in selected_headers:
                    if header.get("text") and header.get("crop_image_path"):
                        selected_headers_dict[header["text"]] = header

                if selected_headers_dict:
                    create_ppt_from_header_dict(selected_headers_dict, pptx_path)
                    st.session_state["search_pptx_path"] = pptx_path
                else:
                    st.session_state["search_pptx_path"] = None
            else:
                st.session_state["search_pptx_path"] = None

            # 다운로드 버튼은 세션 값만 사용
            if st.session_state.get("search_pptx_path") and os.path.exists(
                st.session_state["search_pptx_path"]
            ):
                with open(st.session_state["search_pptx_path"], "rb") as f:
                    st.download_button(
                        label="검색 결과 PPT 다운로드",
                        data=f,
                        file_name="search_results.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    )
            else:
                st.info("검색 결과에 해당하는 헤더 이미지가 없습니다.")
        else:
            st.warning("먼저 벡터스토어에 임베딩을 저장하고, 검색어를 입력하세요.")

with col_upload:
    st.header("📄 파일 업로드")
    uploaded_file = st.file_uploader("PDF 파일을 업로드하세요", type=["pdf"])
    if uploaded_file:
        st.info("파일이 업로드되었습니다. 아래 버튼을 눌러 파싱을 시작하세요.")
        output_dir = "temp_output"
        os.makedirs(output_dir, exist_ok=True)
        md_save_path = os.path.join(
            output_dir, f"temp_{os.path.splitext(uploaded_file.name)[0]}.md"
        )
        pptx_path = os.path.join(output_dir, "exported_slides.pptx")
        progress_bar = st.progress(0, text="대기 중...")
        status_text = st.empty()
        if "parsing_in_progress" not in st.session_state:
            st.session_state["parsing_in_progress"] = False
        if not st.session_state["parsing_in_progress"]:
            if ("pptx_path" in st.session_state and st.session_state["pptx_path"]) or (
                "md_save_path" in st.session_state and st.session_state["md_save_path"]
            ):
                st.success("파싱이 완료되었습니다. 아래에서 파일을 다운로드하세요.")
                pptx_path = st.session_state.get("pptx_path")
                md_save_path = st.session_state.get("md_save_path")
                if pptx_path:
                    try:
                        with open(pptx_path, "rb") as f:
                            st.download_button(
                                label="파싱된 PPT 파일 전체 다운로드",
                                data=f,
                                file_name="exported_slides.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            )
                    except Exception as e:
                        st.error(f"PPT 파일을 찾을 수 없습니다: {e}")
                if md_save_path:
                    try:
                        with open(md_save_path, "rb") as f_md:
                            st.download_button(
                                label="파싱된 Markdown 파일 전체 다운로드",
                                data=f_md,
                                file_name=os.path.basename(md_save_path),
                                mime="text/markdown",
                            )
                    except Exception as e:
                        st.error(f"Markdown 파일을 찾을 수 없습니다: {e}")
            else:
                if st.button("파싱 시작", type="primary"):
                    # temp_output 디렉토리의 모든 파일 삭제
                    for file in os.listdir(output_dir):
                        file_path = os.path.join(output_dir, file)
                        try:
                            if os.path.isfile(file_path):
                                os.remove(file_path)
                        except Exception as e:
                            print(f"Error deleting {file_path}: {e}")

                    st.session_state["parsing_in_progress"] = True
                    st.experimental_rerun()
        if st.session_state["parsing_in_progress"]:
            status_text.info("마크다운 파일로 변환중...")

            temp_pdf_path = os.path.join(output_dir, f"temp_{uploaded_file.name}")
            with open(temp_pdf_path, "wb") as f:
                f.write(uploaded_file.getbuffer())

            # PDF를 이미지로 변환하는 초기 단계 제거
            progress_bar.progress(10, text="문서 파싱 시작...")
            status_text.info("문서 파싱 시작...")

            parser_obj = UpstageParser()
            result = parser_obj.parse(temp_pdf_path)

            # 페이지별로 요소들을 그룹화
            pages = {}
            for element in result.get("elements", []):
                page_num = element.get("page", 1)
                if page_num not in pages:
                    pages[page_num] = []
                pages[page_num].append(element)

            # 헤더 정보 추출 (Upstage coordinates 사용)
            headers = []

            # 각 페이지별로 처리
            for page_num, page_elements in sorted(pages.items()):
                prev_was_heading = False
                consecutive_heading_count = 0  # 연속된 heading의 개수 추적
                first_heading_found = False  # 페이지의 첫 번째 heading 추적
                previous_levels = set()  # 이전에 나온 heading level들을 추적

                for element in page_elements:
                    if element.get("category", "").startswith("heading"):
                        content = element.get("content", {})
                        text = content.get("markdown", "").strip()
                        text = re.sub(r"^#+\s*", "", text)

                        # 한글이 없는 헤더(숫자/영어/특수문자만)는 제외
                        if is_useless_header(text):
                            continue

                        # y 좌표 계산 (상단 두 좌표의 평균)
                        coordinates = element.get("coordinates", [])
                        y_coord = None
                        if coordinates and len(coordinates) >= 4:
                            top_y_coords = [coord["y"] for coord in coordinates[:2]]
                            y_coord = sum(top_y_coords) / len(top_y_coords)

                        # heading level 결정 로직
                        if not first_heading_found:
                            # 페이지의 첫 번째 heading은 항상 level 1
                            final_level = 1
                            first_heading_found = True
                            consecutive_heading_count = 1
                        elif not prev_was_heading:
                            # 본문(heading이 아님) 이후의 heading level 결정
                            if 3 in previous_levels:
                                final_level = 3
                            elif 2 in previous_levels:
                                final_level = 2
                            else:
                                final_level = 2
                            consecutive_heading_count = final_level
                        else:
                            # 이전 요소가 heading이었다면 level 증가
                            consecutive_heading_count += 1
                            final_level = min(
                                consecutive_heading_count, 3
                            )  # 최대 level 3

                        # 현재 level을 previous_levels에 추가
                        previous_levels.add(final_level)

                        # 마크다운 형식의 텍스트 생성
                        markdown_text = "#" * (final_level - 1) + " " + text

                        headers.append(
                            {
                                "page": page_num,
                                "level": final_level,
                                "text": text,
                                "markdown_text": markdown_text,
                                "y": y_coord,
                            }
                        )
                        prev_was_heading = True
                    else:
                        prev_was_heading = False

            # 헤더 딕셔너리 생성
            header_dict = build_header_dictionary(headers, {}, output_dir=output_dir)
            st.session_state["header_dict"] = header_dict

            # 마크다운 문서 생성
            md_docs = []
            current_doc = []

            # 페이지별로 마크다운 문서 생성
            for page_num, page_elements in sorted(pages.items()):
                page_content = []
                for element in page_elements:
                    category = element.get("category", "")
                    content = element.get("content", {})

                    if category.startswith("heading"):
                        # 헤더인 경우
                        text = content.get("markdown", "").strip()
                        # # 기호 제거
                        text = re.sub(r"^#+\s*", "", text)

                        # header_dict에서 해당 헤더의 실제 레벨 사용
                        if text in header_dict:
                            header_info = header_dict[text]
                            level = header_info["level"]
                            page_content.append(f"{'#' * (level)}{" "}{text}")
                            print(f"[헤더 매칭] '{text}' -> Level {level}")
                        else:
                            print(f"[헤더 매칭 실패] '{text}'")
                    else:
                        # 본문인 경우
                        text = content.get("markdown", "").strip()
                        if text:
                            page_content.append(text)

                # 페이지 내용을 하나의 문자열로 합치기
                if page_content:
                    md_docs.append("\n".join(page_content))

            st.session_state["md_docs"] = md_docs
            st.session_state["upstage_parse_result"] = result
            progress_bar.progress(40, text="마크다운 문서 생성 완료!")
            status_text.success("마크다운 문서 생성 완료!")

            # md_docs를 파일로 저장
            md_save_path = os.path.join(output_dir, "parsed_docs.md")
            with open(md_save_path, "w", encoding="utf-8") as f:
                for doc in md_docs:
                    f.write(doc + "\n\n")
            st.session_state["md_save_path"] = md_save_path

            progress_bar.progress(50, text="임베딩 및 벡터스토어 저장중...")
            status_text.info("임베딩 및 벡터스토어 저장중...")

            # 임베딩 및 벡터스토어 생성
            headers_to_split_on = [
                ("##", "Header 2"),  # Header 1 제외
                ("###", "Header 3"),
            ]
            markdown_splitter = MarkdownHeaderTextSplitter(
                headers_to_split_on=headers_to_split_on
            )
            all_chunks = []
            md_docs = st.session_state.get("md_docs", [])

            # chunk 생성 로깅
            print("\n=== Chunk 생성 요약 ===")
            total_chunks = 0
            for page_num, doc in enumerate(md_docs, 1):
                chunks = markdown_splitter.split_text(doc)
                total_chunks += len(chunks)
                print(f"페이지 {page_num}: {len(chunks)}개 chunk")

                # 각 chunk의 헤더 정보만 간단히 출력
                for chunk in chunks:
                    headers = [
                        f"{k}: {v}"
                        for k, v in chunk.metadata.items()
                        if k.startswith("Header")
                    ]
                    if headers:
                        print(f"  - {headers[0]}")
                    chunk.metadata["page_number"] = page_num

                # chunks를 all_chunks에 추가
                all_chunks.extend(chunks)

            print(f"\n총 {total_chunks}개 chunk 생성 완료 (Header 1 제외)")
            print("=====================\n")

            progress_bar.progress(70, text="임베딩 계산중...")
            embeddings = OpenAIEmbeddings()
            db = FAISS.from_documents(all_chunks, embeddings)
            st.session_state["faiss_db"] = db

            progress_bar.progress(80, text="헤더 정보 처리중...")
            status_text.info("헤더 정보 처리중...")

            # 파싱이 끝난 직후 최신 header_dict로 PPT 생성
            create_ppt_from_header_dict(header_dict, pptx_path)
            # header_dict를 표로 시각화
            df = pd.DataFrame(list(header_dict.values()))
            st.subheader("파싱된 주제별 정보 표")
            st.dataframe(df)
            progress_bar.progress(100, text="생성 완료!")
            status_text.success("모든 파일이 파싱된 PPT 생성 완료!")
            st.session_state["pptx_path"] = pptx_path
            st.session_state["parsing_in_progress"] = False
            # 파싱 직후 다운로드 버튼 바로 노출
            st.success("파싱이 완료되었습니다. 아래에서 파일을 다운로드하세요.")
            pptx_path = st.session_state.get("pptx_path")
            md_save_path = st.session_state.get("md_save_path")
            if pptx_path:
                try:
                    with open(pptx_path, "rb") as f:
                        st.download_button(
                            label="파싱된 PPT 파일 다운로드",
                            data=f,
                            file_name="exported_slides.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        )
                except Exception as e:
                    st.error(f"PPT 파일을 찾을 수 없습니다: {e}")
            if md_save_path:
                try:
                    with open(md_save_path, "rb") as f_md:
                        st.download_button(
                            label="파싱된 Markdown 파일 다운로드",
                            data=f_md,
                            file_name=os.path.basename(md_save_path),
                            mime="text/markdown",
                        )
                except Exception as e:
                    st.error(f"Markdown 파일을 찾을 수 없습니다: {e}")
